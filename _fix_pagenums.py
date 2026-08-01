from pptx import Presentation

path = r'_w2_report\proposal.pptx'
prs = Presentation(path)
print(f'Slides: {len(prs.slides)}')

fixed = 0
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '/ 29' in run.text:
                    old = run.text
                    run.text = run.text.replace('/ 29', '/ 28')
                    fixed += 1
                    print(f'  slide {i+1}: "{old}" -> "{run.text}"')

prs.save(path)
print(f'Fixed {fixed} page numbers. Saved.')

# 验证
prs2 = Presentation(path)
remaining = 0
for i, slide in enumerate(prs2.slides):
    for shape in slide.shapes:
        if shape.has_text_frame and '/ 29' in shape.text_frame.text:
            remaining += 1
            print(f'  STILL /29 in slide {i+1}: {shape.text_frame.text[:50]}')
print(f'Remaining /29: {remaining}')
