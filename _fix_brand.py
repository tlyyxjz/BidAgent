# -*- coding: utf-8 -*-
"""将提交材料中的品牌名 BidAgent 改为标小智，保留 GitHub 仓库地址"""
import os
import shutil
import re

SRC = r"C:\Users\Lenovo\Desktop\BidAgent"
DEST = r"C:\Users\Lenovo\Desktop\标小智_GOAI初赛提交"

# === 1. GOAI_初赛提交材料_正式版.md ===
f1 = os.path.join(SRC, "GOAI_初赛提交材料_正式版.md")
with open(f1, "r", encoding="utf-8") as f:
    content = f.read()
original = content

# 品牌名替换（保留 GitHub URL 和文件名）
content = content.replace("BidAgent 是一款", "标小智是一款")
content = content.replace("BidAgent 构建了", "标小智构建了")
# 文件名引用改为标小智
content = content.replace("BidAgent_Demo_脚本.md", "标小智_Demo_脚本.md")

if content != original:
    with open(f1, "w", encoding="utf-8") as f:
        f.write(content)
    print(f"FIXED: {os.path.basename(f1)}")

# === 2. GOAI_初赛提交清单.md ===
f2 = os.path.join(SRC, "_w2_report", "GOAI_初赛提交清单.md")
with open(f2, "r", encoding="utf-8") as f:
    content = f.read()
original = content

# 品牌名替换
content = content.replace("BidAgent_Demo_脚本.md", "标小智_Demo_脚本.md")
content = content.replace("BidAgent_Demo_90s.mp4", "标小智_Demo_90s.mp4")
content = content.replace("BidAgent_第二周任务清单.md", "标小智_第二周任务清单.md")
content = content.replace("封面：BidAgent 智能标讯助手", "封面：标小智")
content = content.replace("项目名称：BidAgent 智能标讯助手", "项目名称：标小智")
content = content.replace("未更新到 BidAgent 六 Agent", "未更新到标小智六 Agent")
content = content.replace("更新到 BidAgent 六 Agent", "更新到标小智六 Agent")

if content != original:
    with open(f2, "w", encoding="utf-8") as f:
        f.write(content)
    print(f"FIXED: {os.path.basename(f2)}")

# === 3. compliance.md ===
f3 = os.path.join(SRC, "_w2_report", "compliance.md")
with open(f3, "r", encoding="utf-8") as f:
    content = f.read()
original = content

content = content.replace("# BidAgent 合规声明", "# 标小智合规声明")
content = content.replace("本文件用于说明 BidAgent 项目", "本文件用于说明标小智项目")
content = content.replace("适用范围：BidAgent 全部", "适用范围：标小智全部")
content = content.replace("完整 BidAgent（证据验证", "完整标小智（证据验证")
content = content.replace("完整 BidAgent（来源", "完整标小智（来源")

if content != original:
    with open(f3, "w", encoding="utf-8") as f:
        f.write(content)
    print(f"FIXED: {os.path.basename(f3)}")

# === 4. PPT - 不含 GitHub URL 的 BidAgent ===
f4 = os.path.join(SRC, "_w2_report", "proposal.pptx")
from pptx import Presentation
prs = Presentation(f4)
ppt_fixed = False
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "BidAgent" in run.text:
                    # 保留 GitHub URL 里的 BidAgent
                    if "github.com" in run.text or "tlyyxjz" in run.text:
                        continue
                    new_text = run.text.replace("BidAgent", "标小智")
                    if new_text != run.text:
                        run.text = new_text
                        ppt_fixed = True
                        print(f"FIXED PPT: {run.text[:60]}")
if ppt_fixed:
    prs.save(f4)
    print(f"FIXED: proposal.pptx")

# === 同步到提交文件夹 ===
print("\n=== 同步到提交文件夹 ===")
sync = [
    (f1, os.path.join(DEST, "01_必交材料", "GOAI_初赛提交材料_正式版.md")),
    (f2, os.path.join(DEST, "01_必交材料", "GOAI_初赛提交清单.md")),
    (f3, os.path.join(DEST, "02_补充材料", "compliance.md")),
    (f4, os.path.join(DEST, "01_必交材料", "proposal.pptx")),
]
for src, dest in sync:
    shutil.copy2(src, dest)
    print(f"SYNCED: {os.path.basename(src)}")

# === 验证 ===
print("\n=== 验证 ===")
for root, dirs, files in os.walk(DEST):
    for fname in files:
        fpath = os.path.join(root, fname)
        rel = os.path.relpath(fpath, DEST)
        if fname.endswith(".pptx"):
            prs = Presentation(fpath)
            for i, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text
                        # 找不含 github 的 BidAgent
                        for line in text.split("\n"):
                            if "BidAgent" in line and "github" not in line.lower():
                                print(f"  STILL: {rel} slide{i+1}: {line[:70]}")
        elif fname.endswith(".md"):
            with open(fpath, "r", encoding="utf-8") as f:
                lines = f.readlines()
            for i, line in enumerate(lines):
                if "BidAgent" in line and "github" not in line.lower() and "tlyyxjz" not in line.lower():
                    print(f"  STILL: {rel} L{i+1}: {line.strip()[:70]}")
print("\n(仅 GitHub URL 中的 BidAgent 保留，其余已改为标小智)")
