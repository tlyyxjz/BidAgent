from pptx import Presentation
pptx_path = r'C:\Users\Lenovo\Desktop\BidAgent\_w2_report\proposal.pptx'
prs = Presentation(pptx_path)

fixed = False
for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '200 项' in run.text:
                    old_text = run.text
                    new_text = run.text.replace('200 项', '826 项')
                    run.text = new_text
                    print(f"SLIDE_{slide_idx+1}: '{old_text}' -> '{new_text}'")
                    fixed = True

if fixed:
    prs.save(pptx_path)
    print("SAVED")
else:
    print("NO_MATCH")
