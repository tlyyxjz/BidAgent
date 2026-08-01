# -*- coding: utf-8 -*-
"""检查提交材料中是否包含敏感数据"""
import os
import re

DEST = r"C:\Users\Lenovo\Desktop\标小智_GOAI初赛提交"

# 敏感数据检查清单 (pattern, 风险等级, 说明)
CHECKS = [
    # 密码/密钥
    (r'@[A-Za-z0-9@!]{8,}!!', "HIGH", "密码模式(含特殊字符)"),
    (r'password\s*[=:]\s*\S+', "HIGH", "password字段"),
    (r'passwd\s*[=:]\s*\S+', "HIGH", "passwd字段"),
    (r'api[_-]?key\s*[=:]\s*\S+', "HIGH", "API key"),
    (r'secret\s*[=:]\s*\S+', "HIGH", "secret字段"),
    (r'token\s*[=:]\s*[A-Za-z0-9]{20,}', "HIGH", "token值"),
    (r'ACCESS_KEY\s*[=:]', "HIGH", "ACCESS_KEY"),
    (r'SECRET_KEY\s*[=:]', "HIGH", "SECRET_KEY"),
    (r'OPENAI_API_KEY\s*[=:]', "HIGH", "OPENAI_API_KEY"),
    (r'DEEPSEEK_API_KEY\s*[=:]', "HIGH", "DEEPSEEK_API_KEY"),
    (r'GLM_API_KEY\s*[=:]', "HIGH", "GLM_API_KEY"),
    (r'sk-[A-Za-z0-9]{20,}', "HIGH", "OpenAI key格式"),
    (r'token_hex\(32\)', "MED", "token_hex调用(代码中允许,文档中不应)"),

    # 手机号
    (r'1[3-9]\d{9}', "HIGH", "手机号"),
    # 邮箱(可能允许,但需检查)
    (r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}', "MED", "邮箱"),

    # 身份证
    (r'\d{17}[\dXx]', "HIGH", "身份证号"),
    # 银行卡
    (r'\d{16,19}', "LOW", "长数字串(可能是银行卡)"),

    # IP 地址(内网)
    (r'192\.168\.\d+\.\d+', "LOW", "内网IP"),
    (r'10\.\d+\.\d+\.\d+', "LOW", "内网IP"),
    (r'172\.(1[6-9]|2\d|3[01])\.\d+\.\d+', "LOW", "内网IP"),

    # 数据库连接串
    (r'mysql://\S+:\S+@', "HIGH", "MySQL连接串"),
    (r'postgresql://\S+:\S+@', "HIGH", "PostgreSQL连接串"),
    (r'sqlite:///.*\.(db|sqlite)', "LOW", "SQLite路径"),

    # 私钥
    (r'-----BEGIN (RSA |EC )?PRIVATE KEY-----', "HIGH", "私钥"),
    (r'-----BEGIN CERTIFICATE-----', "MED", "证书"),

    # 个人信息
    (r'徐浚钊', "LOW", "真实姓名(已知团队成员,可能允许)"),
    (r'王祯明', "LOW", "真实姓名(已知团队成员,可能允许)"),
    (r'上海建桥大学', "LOW", "学校名(已知,可能允许)"),

    # Windows 用户路径(可能泄露用户名)
    (r'C:\\Users\\[^\\"]+', "MED", "Windows用户路径"),
    (r'/Users/[^/"]+', "MED", "macOS用户路径"),
    (r'/home/[^/"]+', "MED", "Linux用户路径"),

    # .env 内容
    (r'\.env\b', "LOW", ".env引用(可能允许)"),

    # GitHub token
    (r'ghp_[A-Za-z0-9]{36}', "HIGH", "GitHub PAT"),
    (r'github_pat_[A-Za-z0-9_]{82}', "HIGH", "GitHub fine-grained PAT"),
]

# 白名单: 这些文件里出现某些内容是允许的
WHITELIST = {
    # 团队成员姓名在所有文件都允许(提交材料本就要署名)
    "徐浚钊": "团队成员",
    "王祯明": "团队成员",
    "上海建桥大学": "团队学校",
    # localhost/127.0.0.1 是 demo 演示地址,允许
    "127.0.0.1": "demo地址",
    "localhost": "demo地址",
    # 邮箱可能用于联系方式
    "13566878907@163.com": "联系邮箱(检查是否允许)",
    "3069969677@qq.com": "邮箱(检查是否允许)",
    # github 用户名
    "tlyyxjz": "github用户名",
    # 比赛名
    "goaihz.com": "比赛官网",
}

print("=" * 80)
print("敏感数据扫描")
print("=" * 80)

issues = []
allowed = []

for root, dirs, files in os.walk(DEST):
    for fname in sorted(files):
        fpath = os.path.join(root, fname)
        rel = os.path.relpath(fpath, DEST)

        # 跳过二进制文件(pptx 等)
        if fname.endswith((".pptx", ".png", ".jpg", ".pdf")):
            if fname.endswith(".pptx"):
                # PPT 用 python-pptx 检查
                try:
                    from pptx import Presentation
                    prs = Presentation(fpath)
                    for i, slide in enumerate(prs.slides):
                        slide_text = ""
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                slide_text += shape.text_frame.text + "\n"
                        for pattern, level, desc in CHECKS:
                            matches = re.findall(pattern, slide_text, re.IGNORECASE)
                            if matches:
                                for m in matches[:3]:  # 每类最多显示3个
                                    m_str = m if isinstance(m, str) else m[0]
                                    # 检查白名单
                                    is_whitelisted = False
                                    for wl_key, wl_desc in WHITELIST.items():
                                        if wl_key in m_str:
                                            allowed.append((rel, f"slide{i+1}", level, desc, m_str[:50], wl_desc))
                                            is_whitelisted = True
                                            break
                                    if not is_whitelisted:
                                        issues.append((rel, f"slide{i+1}", level, desc, m_str[:50]))
                except Exception as e:
                    print(f"  WARN: 无法读取 {fpath}: {e}")
            continue

        # 文本文件
        try:
            with open(fpath, "r", encoding="utf-8") as f:
                content = f.read()
        except:
            try:
                with open(fpath, "r", encoding="gbk") as f:
                    content = f.read()
            except:
                continue

        for pattern, level, desc in CHECKS:
            matches = re.findall(pattern, content, re.IGNORECASE)
            if matches:
                # 去重
                unique = list(set(matches))[:5]
                for m in unique:
                    m_str = m if isinstance(m, str) else m[0]
                    # 检查白名单
                    is_whitelisted = False
                    for wl_key, wl_desc in WHITELIST.items():
                        if wl_key in m_str:
                            allowed.append((rel, "", level, desc, m_str[:60], wl_desc))
                            is_whitelisted = True
                            break
                    if not is_whitelisted:
                        issues.append((rel, "", level, desc, m_str[:60]))

print("-" * 80)
print(f"HIGH 风险: {sum(1 for i in issues if i[2]=='HIGH')}")
print(f"MED  风险: {sum(1 for i in issues if i[2]=='MED')}")
print(f"LOW  风险: {sum(1 for i in issues if i[2]=='LOW')}")
print(f"白名单允许: {len(allowed)}")
print(f"需关注总数: {len(issues)}")
print("-" * 80)

if issues:
    print("\n🔴 需关注的敏感数据:")
    for rel, slide, level, desc, ctx in sorted(issues, key=lambda x: (x[2], x[0])):
        loc = f" {slide}" if slide else ""
        print(f"  [{level}] {rel}{loc}: {desc} -> '{ctx}'")

if allowed:
    print("\n🟢 白名单允许的内容:")
    for rel, slide, level, desc, ctx, reason in sorted(allowed):
        loc = f" {slide}" if slide else ""
        print(f"  [{level}] {rel}{loc}: {desc} -> '{ctx}' ({reason})")
