# -*- coding: utf-8 -*-
"""
修复 BidAgent proposal.pptx 的 slide 13（金融分析页：补充 AHP 权重来源）
与 slide 25（CI/CD 部署页：修正标题与 PORT 信息缺失）。

只修改 proposal.pptx，不动任何 .bak / _backup_ 文件。

用法：
    python _fix_ppt_finance_deploy.py              # 直接覆盖 proposal.pptx
    python _fix_ppt_finance_deploy.py <out.pptx>    # 另存为 <out.pptx>（用于验证/不覆盖原文件）

注意：若 PowerPoint 正打开 proposal.pptx，文件会被独占锁，save 会抛 PermissionError。
请先关闭 PowerPoint（任务管理器结束 POWERPNT，或关闭该演示文稿），再重新运行本脚本。
脚本带幂等保护：若修改已应用，重复运行不会二次追加。
"""
import sys
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPTX = r"C:\Users\Lenovo\Desktop\BidAgent\_w2_report\proposal.pptx"
OUT = sys.argv[1] if len(sys.argv) > 1 else PPTX

GRAY = RGBColor(0x66, 0x66, 0x66)


def find_shape_by_text(slide, needle):
    for shp in slide.shapes:
        if shp.has_text_frame and needle in shp.text_frame.text:
            return shp
    return None


def find_shapes_by_text(slide, needle):
    return [shp for shp in slide.shapes
            if shp.has_text_frame and needle in shp.text_frame.text]


def show(tag, shape):
    if shape is None:
        print(f"  {tag}: <not found>")
        return
    print(f"  {tag}: name={shape.name!r} "
          f"top={round(shape.top/360000,3)}cm h={round(shape.height/360000,3)}cm")
    for pi, p in enumerate(shape.text_frame.paragraphs):
        for ri, r in enumerate(p.runs):
            sz = r.font.size.pt if r.font.size else None
            col = None
            try:
                if r.font.color and r.font.color.rgb is not None:
                    col = str(r.font.color.rgb)
            except Exception:
                col = "<inherit>"
            print(f"    para[{pi}] run[{ri}] = {r.text!r} size={sz} color={col}")


prs = Presentation(PPTX)
print(f"total slides (before): {len(prs.slides)}")

# ---------- Slide 13 (index 12) ----------
slide13 = prs.slides[12]
s13_card = find_shape_by_text(slide13, "集中度 25% + 金额 20%")
print("\n== slide13 供应商信用评分卡片 ==")
show("BEFORE", s13_card)
if s13_card is not None and "AHP" not in s13_card.text_frame.text:
    tf = s13_card.text_frame
    new_para = tf.add_paragraph()
    new_para.alignment = PP_ALIGN.LEFT
    run = new_para.add_run()
    run.text = "权重来源：AHP 层次分析法导出，CR=0.0116<0.1"
    run.font.size = Pt(10)
    run.font.color.rgb = GRAY
    run.font.bold = False
    run.font.name = "Calibri"
    show("AFTER", s13_card)
elif s13_card is not None:
    print("  [skip] AHP 说明已存在，不重复追加")

# ---------- Slide 25 (index 24) ----------
slide25 = prs.slides[24]

print("\n== slide25 标题 ==")
s25_title = find_shape_by_text(slide25, "Docker 一键部署")
show("BEFORE", s25_title)
if s25_title is not None:
    p = s25_title.text_frame.paragraphs[0]
    if p.runs and p.runs[0].text == "Docker 一键部署":
        p.runs[0].text = "Docker 容器化部署 · docker-compose 一键拉起"
        show("AFTER", s25_title)
    else:
        print("  [skip] 标题非原值，已修改过")

print("\n== slide25 web 卡片 ==")
s25_web = find_shape_by_text(slide25, "PORT: 8000")
show("BEFORE", s25_web)
if s25_web is not None and "演示" not in s25_web.text_frame.text:
    tf = s25_web.text_frame
    new_para = tf.add_paragraph()
    new_para.alignment = PP_ALIGN.CENTER
    run = new_para.add_run()
    run.text = "演示：127.0.0.1:8000"
    run.font.size = Pt(9)
    run.font.color.rgb = GRAY
    run.font.bold = False
    run.font.name = "Consolas"
    # 原 0.762cm 偏紧，增高到 1.1cm 以容纳第二行；卡片底 9.144cm，1.1cm 后 9.101cm 仍在卡内
    s25_web.height = Cm(1.1)
    show("AFTER", s25_web)
elif s25_web is not None:
    print("  [skip] 演示地址已存在")

print("\n== slide25 worker / scheduler 卡片 ==")
for idx, shp in enumerate(find_shapes_by_text(slide25, "PORT: -")):
    show(f"BEFORE #{idx}", shp)
    p = shp.text_frame.paragraphs[0]
    if p.runs and p.runs[0].text == "PORT: -":
        p.runs[0].text = "内部服务（不对外暴露）"
        show(f"AFTER #{idx}", shp)
    else:
        print("  [skip] 非原值")

# ---------- 保存 ----------
try:
    prs.save(OUT)
    print(f"\nSaved: {OUT}")
    print(f"total slides (after): {len(prs.slides)}")
except PermissionError as e:
    print(f"\n[SAVE FAILED] PermissionError: {e}")
    print("proposal.pptx 被 PowerPoint 独占锁定。")
    print("请关闭 PowerPoint（任务管理器结束 POWERPNT 进程，或关闭该演示文稿），")
    print("然后重新运行：python _fix_ppt_finance_deploy.py")
    sys.exit(2)
