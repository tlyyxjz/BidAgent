from pptx import Presentation

path = r'_w2_report\proposal.pptx'
prs = Presentation(path)
print(f'Total slides: {len(prs.slides)}')
print()
for i, slide in enumerate(prs.slides):
    title = ''
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            title = shape.text_frame.text.strip().split('\n')[0][:70]
            break
    print(f'Slide {i+1}: {title}')
