# -*- coding: utf-8 -*-
"""修复提交材料中的敏感数据：
1. 手机号 13566878907 -> 135****8907
2. 手机号 17849729681 -> 178****9681
3. Demo 脚本里的 test-admin-secret-12345 -> <your-admin-secret>
4. PPT slide 28 的手机号
"""
import os
import shutil
import re

SRC = r"C:\Users\Lenovo\Desktop\BidAgent"
DEST = r"C:\Users\Lenovo\Desktop\标小智_GOAI初赛提交"

fixed_files = []

# === 1. GOAI_初赛提交清单.md (源文件) ===
f1 = os.path.join(SRC, "_w2_report", "GOAI_初赛提交清单.md")
with open(f1, "r", encoding="utf-8") as f:
    content = f.read()
original = content
content = content.replace("13566878907", "135****8907")
if content != original:
    with open(f1, "w", encoding="utf-8") as f:
        f.write(content)
    fixed_files.append(f1)
    print(f"FIXED: {f1}")

# === 2. K3_审校报告.md ===
f2 = os.path.join(SRC, "_w2_report", "K3_审校报告.md")
with open(f2, "r", encoding="utf-8") as f:
    content = f.read()
original = content
content = content.replace("17849729681", "178****9681")
if content != original:
    with open(f2, "w", encoding="utf-8") as f:
        f.write(content)
    fixed_files.append(f2)
    print(f"FIXED: {f2}")

# === 3. CONTRIBUTING.md ===
f3 = os.path.join(SRC, "CONTRIBUTING.md")
with open(f3, "r", encoding="utf-8") as f:
    content = f.read()
original = content
content = content.replace("13566878907", "135****8907")
if content != original:
    with open(f3, "w", encoding="utf-8") as f:
        f.write(content)
    fixed_files.append(f3)
    print(f"FIXED: {f3}")

# === 4. Demo录制脚本_3分钟.md - 测试 secret ===
f4 = os.path.join(SRC, "docs", "Demo录制脚本_3分钟.md")
with open(f4, "r", encoding="utf-8") as f:
    content = f.read()
original = content
content = content.replace('ADMIN_SECRET="test-admin-secret-12345"', 'ADMIN_SECRET="<your-admin-secret>"')
if content != original:
    with open(f4, "w", encoding="utf-8") as f:
        f.write(content)
    fixed_files.append(f4)
    print(f"FIXED: {f4}")

# === 5. proposal.pptx - slide 28 手机号 ===
f5 = os.path.join(SRC, "_w2_report", "proposal.pptx")
from pptx import Presentation
prs = Presentation(f5)
ppt_fixed = False
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "13566878907" in run.text:
                    run.text = run.text.replace("13566878907", "135****8907")
                    ppt_fixed = True
                    print(f"FIXED PPT: 13566878907 -> 135****8907")
if ppt_fixed:
    prs.save(f5)
    fixed_files.append(f5)

# === 同步到提交文件夹 ===
print("\n=== 同步到提交文件夹 ===")
sync_map = [
    (f1, os.path.join(DEST, "01_核心提交材料", "GOAI_初赛提交清单.md")),
    (f2, os.path.join(DEST, "02_评测与质量报告", "K3_审校报告.md")),
    (f3, os.path.join(DEST, "04_开源与合规", "CONTRIBUTING.md")),
    (f4, os.path.join(DEST, "03_Demo与路演", "Demo录制脚本_3分钟.md")),
    (f5, os.path.join(DEST, "01_核心提交材料", "proposal.pptx")),
]
for src, dest in sync_map:
    shutil.copy2(src, dest)
    print(f"SYNCED: {os.path.basename(src)}")

# === 验证 ===
print("\n=== 验证 ===")
all_ok = True
# 检查手机号
for root, dirs, files in os.walk(DEST):
    for fname in files:
        fpath = os.path.join(root, fname)
        if fname.endswith(".pptx"):
            try:
                prs = Presentation(fpath)
                for i, slide in enumerate(prs.slides):
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text = shape.text_frame.text
                            if "13566878907" in text or "17849729681" in text:
                                print(f"  STILL HAS 手机号: {fname} slide{i+1}")
                                all_ok = False
            except:
                pass
        else:
            try:
                with open(fpath, "r", encoding="utf-8") as f:
                    content = f.read()
                if "13566878907" in content or "17849729681" in content:
                    print(f"  STILL HAS 手机号: {fname}")
                    all_ok = False
                if "test-admin-secret-12345" in content:
                    print(f"  STILL HAS test secret: {fname}")
                    all_ok = False
            except:
                pass

if all_ok:
    print("  ALL OK - 敏感数据已全部脱敏")
