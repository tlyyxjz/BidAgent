"""
将 slide 22 (T7 知识增强规划) 改造为 "LLM 工程化能力" 页
- 标题改为 "02 · LLM 工程化能力"
- 副标题改为突出 LLM 使用 + 成本 + 编排
- 三张卡片改为：意图解析/字段抽取/Agent 编排
- 底部说明改为成本与可复现性
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

path = r'_w2_report\proposal.pptx'
prs = Presentation(path)
slide = prs.slides[21]  # slide 22

shapes = list(slide.shapes)
print(f'Slide 22 has {len(shapes)} shapes')
for j, sh in enumerate(shapes):
    t = sh.text_frame.text.replace('\n', '|')[:90] if sh.has_text_frame else ''
    print(f'  SHAPE_{j}: "{t}"')

# SHAPE_0: 标题
# SHAPE_1: 副标题
# SHAPE_3: 卡片1标题, SHAPE_4: 卡片1正文
# SHAPE_6: 卡片2标题, SHAPE_7: 卡片2正文
# SHAPE_9: 卡片3标题, SHAPE_10: 卡片3正文
# SHAPE_11: 底部说明
# SHAPE_12/13: 页脚/页码

def set_text(shape, lines):
    """设置 text_frame 文本，lines 是字符串列表，每个元素一段"""
    tf = shape.text_frame
    # 清空现有段落
    tf.clear()
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line

# 标题
shapes[0].text_frame.paragraphs[0].runs[0].text = "02 · LLM 工程化能力"

# 副标题
shapes[1].text_frame.paragraphs[0].runs[0].text = "GLM 5.2 驱动意图解析与字段抽取，AgentGraph 串行接力编排，133 万 tokens ¥1.99"

# 卡片1: 意图解析
shapes[3].text_frame.paragraphs[0].runs[0].text = "意图解析 Agent"
set_text(shapes[4], [
    "GLM 5.2 · few-shot + JSON Schema",
    "自然语言 → 5 槽位（关键词/地区/预算/时间/品类）",
    "支持多轮追问：\"只看 100 万以上的\" 增量过滤",
    "失败降级：槽位缺失时主动追问，不静默丢弃",
])

# 卡片2: 字段抽取
shapes[6].text_frame.paragraphs[0].runs[0].text = "字段抽取 Agent"
set_text(shapes[7], [
    "GLM 5.2 · deepseek-v4-flash",
    "99 篇公告 × 6 字段批量抽取",
    "133 万 tokens · ¥1.99 · max_tokens=8000",
    "JSON 输出 + 反幻觉校验 + 无依据字段不展示",
])

# 卡片3: Agent 编排
shapes[9].text_frame.paragraphs[0].runs[0].text = "AgentGraph 编排器"
set_text(shapes[10], [
    "6 Agent 串行接力：意图→采集→加工→质量→金融→交付",
    "节点注册 + 状态传递 + ExecutionTrace 留痕",
    "失败降级：任一 Agent 失败不阻塞整条链路",
    "可审计：每步执行 trace 可回溯",
])

# 底部说明
shapes[11].text_frame.paragraphs[0].runs[0].text = "成本可复现：deepseek-v4-flash · 133 万 tokens · ¥1.99 · 99 篇公告 × 6 字段批量抽取"

prs.save(path)
print('Saved.')

# 验证
prs2 = Presentation(path)
slide2 = prs2.slides[21]
for j, sh in enumerate(slide2.shapes):
    if sh.has_text_frame and sh.text_frame.text.strip():
        print(f'  SHAPE_{j}: "{sh.text_frame.text.replace(chr(10), "|")[:100]}"')
