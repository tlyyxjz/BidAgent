from pptx import Presentation
from copy import deepcopy
from pptx.util import Emu

path = r'_w2_report\proposal.pptx'
prs = Presentation(path)
print(f'Before: {len(prs.slides)} slides')

# 删除 slide 22 (index 21) - T7 知识增强规划
xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)
target_idx = 21
print(f'Target slide {target_idx+1}: ', end='')
target_slide = prs.slides[target_idx]
title = ''
for shape in target_slide.shapes:
    if shape.has_text_frame and shape.text_frame.text:
        title = shape.text_frame.text.split('\n')[0][:60]
        break
print(f'"{title}"')

xml_slides.remove(slides[target_idx])
print(f'After: {len(prs.slides)} slides')

# 更新后续 slide 的页码 N/29 -> N/28 (slide 23-29 的页码)
# 删除后，原 slide 23-29 变成 index 21-27
for new_idx in range(21, len(prs.slides)):
    slide = prs.slides[new_idx]
    old_num = new_idx + 2  # 原 slide 编号 = 新 index + 2 (因为删了一个)
    new_num = new_idx + 1  # 新 slide 编号
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                old_text = f'{old_num} / 29'
                new_text = f'{new_num} / 28'
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)
                    print(f'  slide {new_num}: "{old_text}" -> "{new_text}"')
                # 也处理 29 -> 28 (目录页可能用)
                if '/ 29' in run.text and old_text not in run.text:
                    run.text = run.text.replace('/ 29', '/ 28')
                    print(f'  slide {new_num}: "/ 29" -> "/ 28"')

prs.save(path)
print(f'Saved: {path}')

# 验证
prs2 = Presentation(path)
print(f'Verify: {len(prs2.slides)} slides')
# 检查 T7 是否还在
for i, slide in enumerate(prs2.slides):
    for shape in slide.shapes:
        if shape.has_text_frame and 'T7' in shape.text_frame.text:
            print(f'  WARNING: T7 still in slide {i+1}')
            break
# 检查页码
for i, slide in enumerate(prs2.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text
            if '/ 29' in t:
                print(f'  WARNING: slide {i+1} still has /29: {t[:50]}')
            if '/ 28' in t:
                print(f'  slide {i+1} has /28 page number: OK')
