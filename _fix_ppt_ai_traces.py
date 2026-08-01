# -*- coding: utf-8 -*-
"""Fix AI-generated traces in BidAgent proposal.pptx.

1. Delete all decorative accent line shapes (under titles / card titles).
2. Fix footer / page-number overlap: narrow the full-width footer text box
   "标小智 · GOAI 2026" from 9.0in to 8.0in so it ends at 8.5in, leaving a
   0.5in gap before the page number at 9.0in.
3. Replace old-caliber text:
     20 类 -> 32 类
     三维度加权 -> 五维度加权（25/20/20/15/20）
     活跃度/中标率/偏离度三维度 -> 集中度/金额/频率/地域/采购人五维度
     三维度：活跃度/中标率/偏离度 -> 五维度：集中度/金额/频率/地域/采购人

Does NOT touch proposal.pptx.bak2.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPT = r"C:\Users\Lenovo\Desktop\BidAgent\_w2_report\proposal.pptx"
EMU = 914400
FOOTER_W = int(8.0 * EMU)

# Replacement rules. Applied in order; longest/most-specific first so that
# e.g. "活跃度/中标率/偏离度三维度" is resolved before "三维度加权".
REPLACEMENTS = [
    ("20 类基准价格库", "32 类基准价格库"),
    ("20 类采购品类", "32 类采购品类"),
    ("20 类品类", "32 类品类"),
    ("20 类", "32 类"),  # fallback: guarantees grep "20 类" returns zero
    ("活跃度/中标率/偏离度三维度", "集中度/金额/频率/地域/采购人五维度"),
    ("三维度：活跃度/中标率/偏离度", "五维度：集中度/金额/频率/地域/采购人"),
    ("三维度加权", "五维度加权（25/20/20/15/20）"),
]


def walk(shapes):
    """Yield every shape, descending into GROUP shapes."""
    for shp in shapes:
        yield shp
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                yield from walk(shp.shapes)
            except Exception:
                pass


def is_decorative_line(shp):
    """A thin, wide, text-less LINE or AUTO_SHAPE used as an accent line."""
    st = shp.shape_type
    if st not in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.AUTO_SHAPE):
        return False
    try:
        w = shp.width
        h = shp.height
    except Exception:
        return False
    if w is None or h is None:
        return False
    if not (h < 0.15 * EMU):          # very thin (includes zero-height lines)
        return False
    if not (w >= 0.5 * EMU):          # wide enough to be an accent line
        return False
    txt = ""
    if shp.has_text_frame:
        txt = shp.text_frame.text or ""
    if txt.strip():                   # text boxes are never decorative lines
        return False
    return True


def is_footer_to_fix(shp):
    """The recurring bottom footer '标小智 · GOAI 2026' that overlaps the
    page-number column. Only full-width (>=8.5in) footers need narrowing."""
    if not shp.has_text_frame:
        return False
    txt = shp.text_frame.text or ""
    if not ("标小智" in txt and "GOAI" in txt and "2026" in txt):
        return False
    try:
        w = shp.width
    except Exception:
        return False
    if w is None:
        return False
    return w >= 8.5 * EMU


def apply_replacements_to_text_frame(tf):
    """Apply REPLACEMENTS per paragraph, preserving first-run formatting.
    Tries per-run replacement first (keeps formatting exactly); falls back to
    paragraph-level join when a match spans multiple runs.
    Returns a list of (before, after) for changed paragraphs."""
    changes = []
    for para in tf.paragraphs:
        runs = para.runs
        if not runs:
            continue
        para_old = "".join(r.text for r in runs)

        # 1) per-run replacement (preserves each run's formatting)
        any_run_changed = False
        for r in runs:
            rt = r.text
            new_rt = rt
            for a, b in REPLACEMENTS:
                new_rt = new_rt.replace(a, b)
            if new_rt != rt:
                r.text = new_rt
                any_run_changed = True
        if any_run_changed:
            para_new = "".join(r.text for r in runs)
            if para_new != para_old:
                changes.append((para_old, para_new))
            continue

        # 2) cross-run match: join, replace, put result in first run
        new = para_old
        for a, b in REPLACEMENTS:
            new = new.replace(a, b)
        if new != para_old:
            runs[0].text = new
            for r in runs[1:]:
                r.text = ""
            changes.append((para_old, new))
    return changes


def main():
    prs = Presentation(PPT)

    deleted_lines = 0
    footer_fixed = 0
    text_changes = []  # (slide_idx, before, after)

    for idx, slide in enumerate(prs.slides, start=1):
        # --- 1. collect & delete decorative lines ---
        to_delete = [shp for shp in walk(slide.shapes) if is_decorative_line(shp)]
        for shp in to_delete:
            el = shp._element
            parent = el.getparent()
            if parent is not None:
                parent.remove(el)
                deleted_lines += 1

        # --- 2. narrow full-width footers ---
        for shp in walk(slide.shapes):
            if is_footer_to_fix(shp):
                shp.width = FOOTER_W
                footer_fixed += 1

        # --- 3. replace old-caliber text ---
        for shp in walk(slide.shapes):
            if shp.has_text_frame:
                for old, new in apply_replacements_to_text_frame(shp.text_frame):
                    text_changes.append((idx, old, new))

    prs.save(PPT)

    print("=" * 70)
    print("FIX SUMMARY")
    print("=" * 70)
    print(f"Decorative lines deleted : {deleted_lines}")
    print(f"Footer boxes resized     : {footer_fixed}")
    print(f"Text replacements        : {len(text_changes)}")
    print("-" * 70)
    for idx, old, new in text_changes:
        print(f"  Slide {idx}:")
        print(f"    BEFORE: {old!r}")
        print(f"    AFTER : {new!r}")
    print("=" * 70)


if __name__ == "__main__":
    main()
