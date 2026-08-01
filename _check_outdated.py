# -*- coding: utf-8 -*-
"""精确检查提交文件夹中所有文件的过期口径"""
import os
import re

DEST = r"C:\Users\Lenovo\Desktop\标小智_GOAI初赛提交"

# 关键词检查清单
# (关键词, 是否允许出现, 说明)
CHECKS = [
    ("571", False, "旧测试数，应为826"),
    ("200 项", False, "旧测试数"),
    ("200项", False, "旧测试数"),
    ("三维度", False, "旧口径，应为五维度"),
    ("30/40/30", False, "旧权重，应为25/20/20/15/20"),
    ("T7", False, "已删除的知识增强规划页"),
    ("知识增强规划", False, "已删除的规划页"),
    ("/ 29", False, "旧页码，应为/28"),
    ("/29", False, "旧页码"),
    ("20 类", False, "旧BOQ类数，应为32类"),
    ("20类", False, "旧BOQ类数"),
    ("路线图", False, "已删除的路线图页"),
    ("ScrapeFlow", False, "旧架构名，应为六Agent"),
    ("BidAgent 智能标讯", False, "旧品牌名，应为标小智"),
    ("511", False, "更旧的测试数"),
]

# 允许出现的关键词（历史记录类文档）
HISTORY_FILES = {
    "CHANGELOG.md",  # 变更日志允许记录历史
    "路演讲稿_改动说明.md",  # 改动说明允许提旧值
    "K3_审校报告.md",  # 审校报告记录历史问题
    "K3_交付物复查报告.md",
    "K3_代码review报告.md",
    "K3_测试盲点报告.md",
    "GOAI_初赛提交清单.md",  # 清单记录修复历史
    "credit_score_methodology.md",  # 方法论文档可能提旧维度
}

print("=" * 80)
print("提交材料过期口径精确检查")
print("=" * 80)

issues = []
for root, dirs, files in os.walk(DEST):
    for fname in sorted(files):
        fpath = os.path.join(root, fname)
        rel = os.path.relpath(fpath, DEST)

        # 跳过二进制文件
        if fname.endswith((".pptx", ".png", ".jpg", ".pdf")):
            # PPT 用 python-pptx 检查
            if fname.endswith(".pptx"):
                try:
                    from pptx import Presentation
                    prs = Presentation(fpath)
                    for i, slide in enumerate(prs.slides):
                        for shape in slide.shapes:
                            if not shape.has_text_frame:
                                continue
                            text = shape.text_frame.text
                            for kw, should_be_absent, desc in CHECKS:
                                if kw in text:
                                    is_history = fname in HISTORY_FILES
                                    if not is_history:
                                        issues.append((rel, f"slide{i+1}", kw, desc, text[:60]))
                                        print(f"  ISSUE: {rel} slide{i+1} 含 '{kw}' ({desc})")
                except Exception as e:
                    print(f"  WARN: 无法读取 {fpath}: {e}")
            continue

        # 文本文件
        try:
            with open(fpath, "r", encoding="utf-8") as f:
                content = f.read()
        except:
            try:
                with open(fpath, "r", encoding="gbk") as f:
                    content = f.read()
            except:
                continue

        for kw, should_be_absent, desc in CHECKS:
            if kw in content:
                is_history = fname in HISTORY_FILES
                if is_history:
                    # 历史文档允许，但记录一下
                    count = content.count(kw)
                    print(f"  ALLOWED({count}x): {rel} 含 '{kw}' ({desc}) [历史文档]")
                else:
                    count = content.count(kw)
                    issues.append((rel, "", kw, desc, ""))
                    print(f"  ISSUE({count}x): {rel} 含 '{kw}' ({desc})")

print()
print("=" * 80)
print(f"问题总数: {len(issues)}")
print("=" * 80)
if issues:
    print("\n需要修复的文件:")
    for rel, slide, kw, desc, ctx in issues:
        loc = f" {slide}" if slide else ""
        print(f"  {rel}{loc}: '{kw}' ({desc})")
