"""
P1 价值量化对比：
- slide 5: 右侧 4 个标签框改为"人工 vs 标小智"对比
- slide 14: 底部加一行价值量化总结
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

path = r'_w2_report\proposal.pptx'
prs = Presentation(path)

# === slide 5: 痛点页右侧标签改为量化对比 ===
slide5 = prs.slides[4]
shapes5 = list(slide5.shapes)

# SHAPE_8: "2-4 小时" → "人工 2-4h → 标小智 3min"
# SHAPE_14: "高漏报率" → "人工漏报 40% → 标小智 <5%"
# SHAPE_20: "低效率" → "人工跨 30 平台 → 标小智一键"
# SHAPE_26: "错失机会" → "人工无提醒 → 标小智自动推送"

quant_replacements = {
    8: ("人工 2-4h", "标小智 3min"),
    14: ("人工漏报 40%", "标小智 <5%"),
    20: ("人工跨 30 平台", "标小智一键"),
    26: ("人工无提醒", "标小智自动推"),
}

for shape_idx, (before_text, after_text) in quant_replacements.items():
    shape = shapes5[shape_idx]
    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = before_text
    r1.font.size = Pt(9)
    r1.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    r1.font.bold = False

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "↓"
    r2.font.size = Pt(9)
    r2.font.color.rgb = RGBColor(0x1E, 0x27, 0x61)
    r2.font.bold = True

    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = after_text
    r3.font.size = Pt(10)
    r3.font.color.rgb = RGBColor(0x1E, 0x27, 0x61)
    r3.font.bold = True

    print(f'slide5 SHAPE_{shape_idx}: {before_text} ↓ {after_text}')

# 调整 shape 高度容纳 3 行（原 0.40-0.50 → 0.90）
for shape_idx in quant_replacements:
    shape = shapes5[shape_idx]
    shape.height = Emu(int(0.90 * 914400))
    # top 上移一点保持居中
    old_top = shape.top
    shape.top = old_top - Emu(int(0.20 * 914400))

# === slide 14: Demo 页底部加价值量化总结 ===
slide14 = prs.slides[13]
# 在 6 张卡片下方 (top 6.2) 加一个文本框
from pptx.util import Inches
left = Inches(0.50)
top = Inches(6.30)
width = Inches(9.00)
height = Inches(0.60)
txBox = slide14.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "价值量化：人工跨平台尽调 2-4 小时 → 标小智 3 分钟 · 字段抽取 IoU 0.5339 · 无依据字段率 100%→0% · LLM 成本 ¥1.99/99 篇"
r.font.size = Pt(11)
r.font.color.rgb = RGBColor(0x1E, 0x27, 0x61)
r.font.bold = True

print('slide14: added value quantification footer')

# === slide 14: step 2 意图解析卡片补多轮交互 ===
# SHAPE_12: "展示 5 槽位：关键词/地区/预算/时间/品类"
shapes14 = list(slide14.shapes)
shape12 = shapes14[12]
tf12 = shape12.text_frame
tf12.clear()
p1 = tf12.paragraphs[0]
r1 = p1.add_run()
r1.text = "5 槽位 + 多轮追问"
r1.font.size = Pt(11)
p2 = tf12.add_paragraph()
r2 = p2.add_run()
r2.text = "支持\"只看 100 万以上的\"增量过滤"
r2.font.size = Pt(9)
r2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
print('slide14 SHAPE_12: added multi-turn hint')

prs.save(path)
print('Saved.')

# 验证
prs2 = Presentation(path)
s5 = prs2.slides[4]
for j in [8, 14, 20, 26]:
    sh = list(s5.shapes)[j]
    print(f'  slide5 SHAPE_{j}: "{sh.text_frame.text.replace(chr(10), "|")}"')
s14 = prs2.slides[13]
print(f'  slide14 last shapes:')
for j, sh in enumerate(list(s14.shapes)[-3:]):
    if sh.has_text_frame:
        print(f'    SHAPE_{len(list(s14.shapes))-3+j}: "{sh.text_frame.text.replace(chr(10), "|")[:100]}"')
