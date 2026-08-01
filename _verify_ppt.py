from pptx import Presentation
prs = Presentation(r'C:\Users\Lenovo\Desktop\BidAgent\_w2_report\proposal.pptx')
print(f"SLIDE_COUNT: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        texts.append(run.text.strip())
    full = " | ".join(texts)
    print(f"SLIDE_{i+1}: {full[:300]}")
