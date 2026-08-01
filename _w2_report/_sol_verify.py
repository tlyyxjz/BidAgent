# -*- coding: utf-8 -*-
"""Sol 严格规则 PPT 复核脚本 - 实际运行 python-pptx 读取验证"""
import sys
import io
from pptx import Presentation
from pptx.util import Emu, Inches

# 强制 stdout utf-8 (Windows 控制台兼容)
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

PPT_PATH = r"C:\Users\Lenovo\Desktop\BidAgent\_w2_report\proposal.pptx"

prs = Presentation(PPT_PATH)
SLIDE_W_IN = prs.slide_width / 914400.0
SLIDE_H_IN = prs.slide_height / 914400.0
print(f"[INFO] 共 {len(prs.slides)} 张幻灯片, slide 宽={SLIDE_W_IN:.3f}in 高={SLIDE_H_IN:.3f}in")
print("=" * 80)


def emu_to_in(v):
    if v is None:
        return None
    return v / 914400.0


def shape_text(shape):
    """递归提取 shape 的全部文本"""
    parts = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text:
                    parts.append(run.text)
            # 段落结尾也算
            if para.text:
                # 已经在 runs 里
                pass
    # 表格
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                parts.append(cell.text)
    # 组合
    if shape.shape_type == 6:  # GROUP
        for sub in shape.shapes:
            parts.append(shape_text(sub))
    return " | ".join(parts)


def all_shapes(slide):
    """扁平化返回所有 shape（含组合内的子 shape）, 返回 [(shape, parent_path)]"""
    result = []

    def _walk(shapes, path=""):
        for sh in shapes:
            result.append((sh, path))
            if sh.shape_type == 6:  # GROUP
                _walk(sh.shapes, path + "/group")
    _walk(slide.shapes)
    return result


def slide_full_text(slide):
    parts = []
    for sh, _ in all_shapes(slide):
        t = shape_text(sh)
        if t:
            parts.append(t)
    return "\n".join(parts)


# ============================================================
# 核查 1：5 项修改是否真实落地
# ============================================================
print("\n【核查 1：5 项修改是否真实落地】")
print("-" * 80)

slides = list(prs.slides)

# --- 1) slide 22 (1-based) ---
print("\n[1.1] slide 22 应为 '02 · LLM 工程化能力' 页")
if len(slides) >= 22:
    s22 = slides[21]  # 0-based
    txt22 = slide_full_text(s22)
    print(f"  slide 22 全文前 600 字: {txt22[:600]}")
    checks = [
        ("02 · LLM 工程化能力", "02 · LLM 工程化能力" in txt22 or "LLM 工程化" in txt22),
        ("不是 T7 知识增强规划", "T7" not in txt22 and "知识增强规划" not in txt22),
        ("卡片: 意图解析 Agent", "意图解析" in txt22 and "Agent" in txt22),
        ("卡片: 字段抽取 Agent", "字段抽取" in txt22),
        ("卡片: AgentGraph 编排器", "AgentGraph" in txt22),
        ("底部: 133 万 tokens", "133 万" in txt22 or "133万" in txt22),
        ("底部: ¥1.99", "1.99" in txt22),
    ]
    for label, ok in checks:
        print(f"  {'✅' if ok else '❌'} {label}")
else:
    print("  ❌ PPT 不足 22 页")

# --- 2) slide 5 痛点页右侧 4 个标签框 ---
print("\n[1.2] slide 5 痛点页右侧 4 个标签框应含对比文本")
if len(slides) >= 5:
    s5 = slides[4]
    txt5 = slide_full_text(s5)
    print(f"  slide 5 全文前 800 字: {txt5[:800]}")
    checks = [
        ("人工 2-4h", "2-4" in txt5 and "人工" in txt5),
        ("标小智 3min", "3min" in txt5 or "3 min" in txt5 or "3分钟" in txt5),
    ]
    for label, ok in checks:
        print(f"  {'✅' if ok else '❌'} {label}")
else:
    print("  ❌ PPT 不足 5 页")

# --- 3) slide 14 Demo 页 ---
print("\n[1.3] slide 14 Demo 页底部价值量化文本")
if len(slides) >= 14:
    s14 = slides[13]
    txt14 = slide_full_text(s14)
    print(f"  slide 14 全文前 800 字: {txt14[:800]}")
    checks = [
        ("价值量化 / 2-4 小时", "2-4" in txt14 and ("小时" in txt14 or "h" in txt14)),
        ("3 分钟", "3 分钟" in txt14 or "3分钟" in txt14),
    ]
    for label, ok in checks:
        print(f"  {'✅' if ok else '❌'} {label}")
else:
    print("  ❌ PPT 不足 14 页")

# --- 4) slide 14 step 2 意图解析卡片含"多轮追问" ---
print("\n[1.4] slide 14 step 2 意图解析卡片含 '多轮追问'")
if len(slides) >= 14:
    has_dz = "多轮追问" in txt14
    print(f"  {'✅' if has_dz else '❌'} 多轮追问 在 slide 14: {has_dz}")
else:
    print("  ❌ PPT 不足 14 页")

# --- 5) slide 26 开放复用页 ---
print("\n[1.5] slide 26 开放复用页 github 地址")
if len(slides) >= 26:
    s26 = slides[25]
    txt26 = slide_full_text(s26)
    print(f"  slide 26 全文前 800 字: {txt26[:800]}")
    checks = [
        ("github 地址", "github.com/tlyyxjz/BidAgent" in txt26 or "tlyyxjz/BidAgent" in txt26),
        ("欢迎 issue/PR", "issue" in txt26 and "PR" in txt26),
    ]
    for label, ok in checks:
        print(f"  {'✅' if ok else '❌'} {label}")
else:
    print("  ❌ PPT 不足 26 页")

print("\n" + "=" * 80)
print("【核查 2：布局重叠/溢出检查】")
print("-" * 80)

TOL = 0.05  # 0.05 inch 容差


def fmt_in(v):
    if v is None:
        return "None"
    return f"{v:.3f}"


for idx, slide in enumerate(slides, start=1):
    overflow_w = []
    overflow_h = []
    overlaps = []
    shape_infos = []  # (name, left, top, width, height, right, bottom)
    for sh, path in all_shapes(slide):
        try:
            l = emu_to_in(sh.left)
            t = emu_to_in(sh.top)
            w = emu_to_in(sh.width)
            h = emu_to_in(sh.height)
        except Exception:
            continue
        if l is None or t is None or w is None or h is None:
            continue
        # 跳过 None 占位 (一些 shape 默认值)
        if w <= 0 or h <= 0:
            continue
        right = l + w
        bottom = t + h
        name = sh.name or sh.shape_type.__str__()
        shape_infos.append((name, l, t, w, h, right, bottom))
        if right > SLIDE_W_IN + TOL:
            overflow_w.append((name, l, w, right))
        if bottom > SLIDE_H_IN + TOL:
            overflow_h.append((name, t, h, bottom))

    # 重叠检测：仅对非背景 shape 报告显著重叠
    # 启发式：跳过 name 含 "背景/background/rect" 等
    for i in range(len(shape_infos)):
        for j in range(i + 1, len(shape_infos)):
            n1, l1, t1, w1, h1, r1, b1 = shape_infos[i]
            n2, l2, t2, w2, h2, r2, b2 = shape_infos[j]
            # 计算 IOU
            ix = max(0, min(r1, r2) - max(l1, l2))
            iy = max(0, min(b1, b2) - max(t1, t2))
            if ix <= 0 or iy <= 0:
                continue
            inter = ix * iy
            area1 = w1 * h1
            area2 = w2 * h2
            union = area1 + area2 - inter
            iou = inter / union if union > 0 else 0
            # 显著重叠：IOU > 0.20 且交集面积较大
            if iou > 0.20 and inter > 0.5:
                # 排除明显的背景卡片（width/height 大于 8in）
                if max(w1, h1) > 8 or max(w2, h2) > 8:
                    continue
                overlaps.append((n1, n2, iou, inter))

    if overflow_w or overflow_h or overlaps:
        print(f"\n  slide {idx}:")
        for n, l, w, r in overflow_w:
            print(f"    ⚠️ 宽度溢出: {n} left={fmt_in(l)} width={fmt_in(w)} right={fmt_in(r)} (slide_w={SLIDE_W_IN:.3f})")
        for n, t, h, b in overflow_h:
            print(f"    ⚠️ 高度溢出: {n} top={fmt_in(t)} height={fmt_in(h)} bottom={fmt_in(b)} (slide_h={SLIDE_H_IN:.3f})")
        for n1, n2, iou, inter in overlaps:
            print(f"    ⚠️ 重叠: {n1} <-> {n2} IOU={iou:.2f} 交集={inter:.2f}sq.in")

# 重点关注 slide 5/14/22/26
print("\n  --- 重点 slide 形状明细 ---")
for tag, idx in [("slide 5", 5), ("slide 14", 14), ("slide 22", 22), ("slide 26", 26)]:
    if len(slides) >= idx:
        s = slides[idx - 1]
        print(f"\n  {tag} 形状清单:")
        for sh, path in all_shapes(s):
            try:
                l = emu_to_in(sh.left)
                t = emu_to_in(sh.top)
                w = emu_to_in(sh.width)
                h = emu_to_in(sh.height)
            except Exception:
                continue
            if l is None or w is None or w <= 0:
                continue
            txt = shape_text(sh)[:60].replace("\n", " ")
            print(f"    {sh.name}: L={fmt_in(l)} T={fmt_in(t)} W={fmt_in(w)} H={fmt_in(h)} R={fmt_in(l + w) if w else '-'} B={fmt_in(t + h) if h else '-'} | {txt}")

print("\n" + "=" * 80)
print("【核查 3：口径一致性】")
print("-" * 80)

KEYWORDS = [
    ("20 类", "应零匹配(改 32 类)"),
    ("三维度", "应零匹配(改 五维度)"),
    ("30/40/30", "应零匹配"),
    ("T7", "应零匹配(已删除/改造)"),
    ("知识增强规划", "应零匹配"),
    ("/ 29", "应零匹配(改 / 28)"),
    ("/29", "应零匹配(改 /28)"),
    ("571", "应零匹配(改 826)"),
]
for kw, note in KEYWORDS:
    hits = []
    for i, s in enumerate(slides, 1):
        t = slide_full_text(s)
        if kw in t:
            # 找出具体在哪个 shape
            for sh, _ in all_shapes(s):
                st = shape_text(sh)
                if kw in st:
                    hits.append((i, sh.name, st[:80].replace("\n", " ")))
    status = "✅" if not hits else "❌"
    print(f"\n  {status} 关键词 '{kw}' ({note}): 命中 {len(hits)} 处")
    for i, n, snippet in hits[:10]:
        print(f"     - slide {i} shape={n}: {snippet}")

print("\n" + "=" * 80)
print("【核查 4：页码正确性】")
print("-" * 80)

import re
page_pat = re.compile(r'(\d+)\s*/\s*28')
for i, s in enumerate(slides, 1):
    txt = slide_full_text(s)
    matches = page_pat.findall(txt)
    # 也找 / 29
    bad_pat = re.compile(r'(\d+)\s*/\s*(\d+)')
    all_matches = bad_pat.findall(txt)
    if matches:
        # 期望 N == i
        nums = [int(m) for m in matches]
        ok = nums[0] == i
        print(f"  slide {i}: 页码 = {matches} {'✅' if ok else '❌'}")
    else:
        # 是否有任何 x/y 形式
        if all_matches:
            print(f"  slide {i}: ❌ 页码格式异常 {all_matches} (非 /28)")
        else:
            print(f"  slide {i}: ⚠️ 未找到 N/28 形式页码")

print("\n" + "=" * 80)
print("【核查 5：AI 常见错误复查】")
print("-" * 80)

# 5.1 标题下装饰线 - 检查每张 slide 是否有紧贴标题下方的细矩形
print("\n[5.1] 标题下装饰线检测（细矩形 shape 紧贴标题下方）")
for i, s in enumerate(slides, 1):
    title_shape = None
    line_shapes = []
    for sh, _ in all_shapes(s):
        # 找标题（最顶部的文本框，且较大字体）
        if sh.has_text_frame and sh.left is not None and sh.top is not None:
            t = emu_to_in(sh.top)
            try:
                txt = sh.text_frame.text.strip()
            except Exception:
                txt = ""
            # 简单判断：标题一般在 top < 1.5 且有文本
            if t is not None and t < 1.5 and txt:
                # 检查字体内是否较大
                max_font = 0
                try:
                    for para in sh.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                pt = run.font.size.pt
                                if pt > max_font:
                                    max_font = pt
                except Exception:
                    pass
                if max_font >= 24 or (title_shape is None):
                    if max_font >= 24:
                        title_shape = (sh, t, emu_to_in(sh.height), max_font)
        # 找窄长的矩形（装饰线特征：width >> height，且 height < 0.15）
        try:
            w = emu_to_in(sh.width)
            h = emu_to_in(sh.height)
            t = emu_to_in(sh.top)
            if w and h and w > 1.5 and h < 0.15 and t is not None:
                line_shapes.append((sh.name, t, w, h))
        except Exception:
            pass
    if title_shape and line_shapes:
        ts, tt, th, tfont = title_shape
        # 装饰线紧贴标题下方：top 在 [title_bottom - 0.1, title_bottom + 0.5]
        tb = (tt or 0) + (th or 0)
        for n, lt, lw, lh in line_shapes:
            if lt is not None and tb - 0.1 <= lt <= tb + 0.6:
                print(f"  ⚠️ slide {i}: 标题(top={tt:.2f},font={tfont}) 下方疑似装饰线 shape={n} top={lt:.2f} w={lw:.2f} h={lh:.2f}")

# 5.2 占位符残留
print("\n[5.2] 占位符残留检查")
PLACEHOLDERS = ["xxx", "XXX", "lorem", "ipsum", "占位", "待填写", "待补充", "TBD", "TODO", "此处", "your text", "请输入"]
for i, s in enumerate(slides, 1):
    txt = slide_full_text(s)
    low = txt.lower()
    for ph in PLACEHOLDERS:
        if ph.lower() in low:
            # 找出具体 shape
            for sh, _ in all_shapes(s):
                st = shape_text(sh)
                if ph.lower() in st.lower():
                    print(f"  ⚠️ slide {i}: 占位符 '{ph}' 出现在 shape={sh.name}: {st[:80]}")

# 5.3 "200 项" 测试数残留
print("\n[5.3] '200 项' 测试数残留检查（应改 826）")
for i, s in enumerate(slides, 1):
    txt = slide_full_text(s)
    if "200 项" in txt or "200项" in txt:
        print(f"  ❌ slide {i}: 仍含 '200 项'")
    # 顺便找 200 后跟 项/测 之类
    import re as _re
    if _re.search(r"200\s*项", txt):
        print(f"  ❌ slide {i}: 正则匹配 '200\\s*项'")

# 5.4 "本地" 字样在部署页
print("\n[5.4] '本地' 字样在部署页检查")
# 部署页通常含 "部署" 关键词，先定位
deploy_slides = []
for i, s in enumerate(slides, 1):
    txt = slide_full_text(s)
    if "部署" in txt:
        deploy_slides.append((i, txt))
for i, txt in deploy_slides:
    if "本地" in txt:
        # 找具体 shape
        for sh, _ in all_shapes(slides[i - 1]):
            st = shape_text(sh)
            if "本地" in st:
                print(f"  ⚠️ slide {i}(部署页): 含 '本地' 在 shape={sh.name}: {st[:80]}")
    else:
        print(f"  ℹ️ slide {i}(部署页): 未发现 '本地' 字样")

# 5.5 路线图/未来规划页残留
print("\n[5.5] 路线图/未来规划页残留检查")
ROADMAP_KW = ["路线图", "未来规划", "Roadmap", "roadmap", "未来工作", "下一步规划", "规划路线"]
for i, s in enumerate(slides, 1):
    txt = slide_full_text(s)
    for kw in ROADMAP_KW:
        if kw in txt:
            for sh, _ in all_shapes(s):
                st = shape_text(sh)
                if kw in st:
                    print(f"  ⚠️ slide {i}: 含 '{kw}' 在 shape={sh.name}: {st[:80]}")

print("\n" + "=" * 80)
print("[DONE] 复核脚本执行完毕")
