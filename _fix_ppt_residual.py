"""修复 BidAgent proposal.pptx 中残留的 3 处旧口径问题。

旧口径（三维度 30/40/30）→ 新口径（五维度 25/20/20/15/20）

替换规则：
1. "活跃度 30% + 中标率 40% + 偏离度 30%" → "集中度 25% + 金额 20% + 频率 20% + 地域 15% + 采购人 20%"
2. "活跃度30%+中标率40%+偏离度30%"（无空格版）→ 同上新口径
3. "30/40/30" 仅在权重语境下（同段落或紧邻 shape 含 "加权/权重/评分/三维度/五维度"）→ "25/20/20/15/20"

不动备份文件（.bak / .bak2 / _backup_）。
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

PPT_PATH = Path(r"C:\Users\Lenovo\Desktop\BidAgent\_w2_report\proposal.pptx")

OLD_FULL = "活跃度 30% + 中标率 40% + 偏离度 30%"
OLD_FULL_NOSPACE = "活跃度30%+中标率40%+偏离度30%"
NEW_FULL = "集中度 25% + 金额 20% + 频率 20% + 地域 15% + 采购人 20%"

OLD_RATIO = "30/40/30"
NEW_RATIO = "25/20/20/15/20"

# 权重语境关键词
WEIGHT_KEYWORDS = ("加权", "权重", "评分", "三维度", "五维度")


def _shape_has_weight_context(shp) -> bool:
    """检查一个 shape 自身的文本是否含权重语境关键词。"""
    if not shp.has_text_frame:
        return False
    txt = shp.text_frame.text
    return any(kw in txt for kw in WEIGHT_KEYWORDS)


def _slide_has_weight_context(slide, exclude_shp) -> bool:
    """检查整张 slide 上（排除某 shape）是否有权重语境关键词。"""
    for shp in slide.shapes:
        if shp is exclude_shp:
            continue
        if _shape_has_weight_context(shp):
            return True
    return False


def main() -> int:
    if not PPT_PATH.exists():
        print(f"[ERROR] PPT not found: {PPT_PATH}", file=sys.stderr)
        return 2

    prs = Presentation(str(PPT_PATH))
    slide_total = len(prs.slides)
    print(f"[INFO] slides total = {slide_total}")

    replacements: list[tuple[int, str, int, str, str]] = []
    # (slide_idx, shape_name, para_idx, before, after)

    for idx, slide in enumerate(prs.slides, start=1):
        for shp in slide.shapes:
            if not shp.has_text_frame:
                continue
            tf = shp.text_frame
            for p_idx, para in enumerate(tf.paragraphs):
                para_text = "".join(r.text for r in para.runs)
                if not para_text:
                    continue

                # 规则 1 & 2：完整口径串替换（含无空格版本）
                if OLD_FULL in para_text or OLD_FULL_NOSPACE in para_text:
                    new_para_text = para_text.replace(OLD_FULL, NEW_FULL).replace(
                        OLD_FULL_NOSPACE, NEW_FULL
                    )
                    runs = para.runs
                    if runs:
                        runs[0].text = new_para_text
                        for r in runs[1:]:
                            r.text = ""
                    replacements.append((idx, shp.name, p_idx, para_text, new_para_text))
                    continue

                # 规则 3：30/40/30 仅在权重语境下替换
                if OLD_RATIO in para_text:
                    in_self = _shape_has_weight_context(shp)
                    in_slide = _slide_has_weight_context(slide, shp)
                    if in_self or in_slide:
                        new_para_text = para_text.replace(OLD_RATIO, NEW_RATIO)
                        runs = para.runs
                        if runs:
                            runs[0].text = new_para_text
                            for r in runs[1:]:
                                r.text = ""
                        replacements.append((idx, shp.name, p_idx, para_text, new_para_text))
                    else:
                        print(
                            f"[WARN] slide {idx} shape='{shp.name}' para#{p_idx} "
                            f"contains '{OLD_RATIO}' but NOT in weight context, skipped: {para_text!r}"
                        )

    if not replacements:
        print("[WARN] No replacements made.")
        return 0

    print("\n[INFO] Replacements to apply:")
    for slide_idx, shape_name, p_idx, before, after in replacements:
        print(f"  - slide {slide_idx} / shape='{shape_name}' / para#{p_idx}")
        print(f"      before: {before!r}")
        print(f"      after : {after!r}")

    prs.save(str(PPT_PATH))
    print(f"\n[OK] Saved: {PPT_PATH}")

    # ---- 验证 ----
    print("\n[VERIFY] Re-opening PPT for verification...")
    prs2 = Presentation(str(PPT_PATH))
    new_total = len(prs2.slides)
    print(f"[VERIFY] slides total = {new_total}")

    hits_old_full = []
    hits_old_ratio_in_weight = []
    hits_new_full = []
    hits_new_ratio = []

    for idx, slide in enumerate(prs2.slides, start=1):
        for shp in slide.shapes:
            if not shp.has_text_frame:
                continue
            txt = shp.text_frame.text
            if OLD_FULL in txt or OLD_FULL_NOSPACE in txt:
                hits_old_full.append((idx, shp.name, txt))
            if OLD_RATIO in txt:
                if _shape_has_weight_context(shp) or _slide_has_weight_context(slide, shp):
                    hits_old_ratio_in_weight.append((idx, shp.name, txt))
            if NEW_FULL in txt:
                hits_new_full.append((idx, shp.name, txt))
            if NEW_RATIO in txt:
                hits_new_ratio.append((idx, shp.name, txt))

    print(f"[VERIFY] grep '活跃度 30% / 活跃度30%' matches = {len(hits_old_full)}")
    for h in hits_old_full:
        print(f"    !! slide {h[0]} shape='{h[1]}': {h[2]!r}")
    print(
        f"[VERIFY] grep '30/40/30' in weight context matches = "
        f"{len(hits_old_ratio_in_weight)}"
    )
    for h in hits_old_ratio_in_weight:
        print(f"    !! slide {h[0]} shape='{h[1]}': {h[2]!r}")
    print(f"[VERIFY] grep new '集中度 25% + ...' matches = {len(hits_new_full)}")
    for h in hits_new_full:
        print(f"    ok slide {h[0]} shape='{h[1]}': {h[2]!r}")
    print(f"[VERIFY] grep new '25/20/20/15/20' matches = {len(hits_new_ratio)}")
    for h in hits_new_ratio:
        print(f"    ok slide {h[0]} shape='{h[1]}': {h[2]!r}")

    print(f"\n[VERIFY] slide total preserved (29) = {new_total == 29}")

    ok = (
        not hits_old_full
        and not hits_old_ratio_in_weight
        and hits_new_full
        and hits_new_ratio
        and new_total == 29
    )
    print(f"\n[RESULT] {'PASS' if ok else 'FAIL'}")
    return 0 if ok else 1


if __name__ == "__main__":
    sys.exit(main())
