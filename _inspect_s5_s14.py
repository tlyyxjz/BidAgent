from pptx import Presentation

path = r'_w2_report\proposal.pptx'
prs = Presentation(path)

for idx in [4, 13]:  # slide 5, slide 14
    slide = prs.slides[idx]
    print(f'\n=== SLIDE {idx+1} ===')
    for j, shape in enumerate(slide.shapes):
        left = shape.left/914400 if shape.left else 0
        top = shape.top/914400 if shape.top else 0
        w = shape.width/914400 if shape.width else 0
        h = shape.height/914400 if shape.height else 0
        text = ''
        if shape.has_text_frame:
            text = shape.text_frame.text.replace('\n', '|')[:120]
        print(f'  SHAPE_{j} type={shape.shape_type} pos=({left:.2f},{top:.2f}) size=({w:.2f}x{h:.2f}) text="{text}"')
