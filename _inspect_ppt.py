# -*- coding: utf-8 -*-
"""
BidAgent proposal.pptx 视觉审查脚本 (refined)
- 提取每张 slide 的所有 shape 位置和尺寸
- 自动检测 AI 做幻灯片常见错误：
  a. 元素溢出 (overflow)
  b. 元素重叠 (overlap) — 过滤容器-子元素包含关系，减少误报
  c. 边距不足 (margin)
  d. 文本框过窄且含长文本
  e. 标题装饰线 (LINE shape 位于标题下方 -> AI 生成痕迹)
输出: C:\\Users\\Lenovo\\Desktop\\BidAgent\\_ppt_inspection_report.txt
"""
from __future__ import annotations
import os
import sys
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPT_PATH = r"C:\Users\Lenovo\Desktop\BidAgent\_w2_report\proposal.pptx"
REPORT_PATH = r"C:\Users\Lenovo\Desktop\BidAgent\_ppt_inspection_report.txt"

EMU_PER_INCH = 914400

# 阈值
MARGIN_INCH = 0.3            # 边距不足阈值
NARROW_WIDTH_INCH = 2.0      # 文本框过窄阈值
LONG_TEXT_LEN = 25           # 长文本判定字符数
OVERLAP_TOL_INCH = 0.05      # 容差，避免贴边误报


def is_background_or_decoration(shape, slide_w, slide_h):
    try:
        if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
            return True  # 无定位信息，跳过
        w_in = shape.width / EMU_PER_INCH
        h_in = shape.height / EMU_PER_INCH
        # 占据 >= 85% 宽度且 >= 85% 高度 -> 视为背景
        if w_in >= (slide_w / EMU_PER_INCH) * 0.85 and h_in >= (slide_h / EMU_PER_INCH) * 0.85:
            return True
    except Exception:
        return True
    return False


def rect_intersect(a, b, tol=OVERLAP_TOL_INCH * EMU_PER_INCH):
    if any(v is None for v in a + b):
        return False
    al, at, ar, ab = a
    bl, bt, br, bb = b
    inter_l = max(al, bl)
    inter_t = max(at, bt)
    inter_r = min(ar, br)
    inter_b = min(ab, bb)
    if inter_r - inter_l <= tol or inter_b - inter_t <= tol:
        return False
    return True


def rect_contains(outer, inner, tol=OVERLAP_TOL_INCH * EMU_PER_INCH):
    """outer 是否完全包含 inner (容器-子元素关系)"""
    if any(v is None for v in outer + inner):
        return False
    ol, ot, orr, ob = outer
    il, it, ir, ib = inner
    return (il >= ol - tol) and (it >= ot - tol) and (ir <= orr + tol) and (ib <= ob + tol)


def shape_rect(shape):
    if None in (shape.left, shape.top, shape.width, shape.height):
        return None
    l = shape.left
    t = shape.top
    return (l, t, l + shape.width, t + shape.height)


def get_shape_type_name(shape):
    try:
        return str(shape.shape_type)
    except Exception:
        return "UNKNOWN"


def is_line_shape(shape):
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            return True
    except Exception:
        pass
    try:
        el = shape._element
        if el.tag.endswith("}cxnSp"):
            return True
    except Exception:
        pass
    # 极扁的 AUTO_SHAPE 也视作装饰线 (height < 0.12in 且 width >= 0.3in, 排除小圆点)
    try:
        if shape.height is not None and shape.width is not None:
            h_in = shape.height / EMU_PER_INCH
            w_in = shape.width / EMU_PER_INCH
            if h_in < 0.12 and w_in >= 0.3:
                return True
    except Exception:
        pass
    return False


def get_text(shape):
    try:
        if shape.has_text_frame:
            return shape.text_frame.text
    except Exception:
        pass
    return ""


def main():
    if not os.path.exists(PPT_PATH):
        print(f"ERROR: PPT not found: {PPT_PATH}")
        sys.exit(1)

    prs = Presentation(PPT_PATH)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    slide_count = len(prs.slides)

    lines_out = []
    lines_out.append("=" * 70)
    lines_out.append("BidAgent proposal.pptx 视觉审查报告")
    lines_out.append("=" * 70)
    lines_out.append(f"文件: {PPT_PATH}")
    lines_out.append(f"SLIDE_SIZE: {slide_w/EMU_PER_INCH:.2f} x {slide_h/EMU_PER_INCH:.2f} inches")
    lines_out.append(f"SLIDE_SIZE_EMU: {slide_w} x {slide_h}")
    lines_out.append(f"SLIDE_COUNT: {slide_count}")
    lines_out.append(f"检测阈值: margin<{MARGIN_INCH}in, narrow_text<{NARROW_WIDTH_INCH}in (text>={LONG_TEXT_LEN} chars), overlap_tol={OVERLAP_TOL_INCH}in")
    lines_out.append("说明: 已过滤容器-子元素包含关系; 极扁的 AUTO_SHAPE (h<0.12in,w>=1in) 视作装饰线")
    lines_out.append("")

    severe_issues_total = 0
    warnings_total = 0
    ok_total = 0

    # 收集跨 slide 汇总
    slide_summaries = []

    for i, slide in enumerate(prs.slides):
        slide_idx = i + 1
        lines_out.append(f"\n=== SLIDE {slide_idx} ===")

        shapes_info = []
        for j, shape in enumerate(slide.shapes):
            try:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                text = get_text(shape)
                type_name = get_shape_type_name(shape)
                is_bg = is_background_or_decoration(shape, slide_w, slide_h)
                is_ln = is_line_shape(shape)
                rect = shape_rect(shape)
                right = (left + width) if (left is not None and width is not None) else None
                bottom = (top + height) if (top is not None and height is not None) else None
                shapes_info.append({
                    "idx": j, "type": type_name, "is_line": is_ln,
                    "left": left, "top": top, "width": width, "height": height,
                    "right": right, "bottom": bottom, "rect": rect,
                    "text": text, "is_bg": is_bg,
                })
            except Exception as e:
                shapes_info.append({
                    "idx": j, "type": "ERROR", "is_line": False,
                    "left": None, "top": None, "width": None, "height": None,
                    "right": None, "bottom": None, "rect": None,
                    "text": "", "is_bg": True, "error": str(e),
                })
                lines_out.append(f"  [ERROR] SHAPE_{j}: 读取失败 - {e}")

        slide_warnings = []
        slide_severe = []

        # a. 元素溢出
        for s in shapes_info:
            if s["left"] is None or s["top"] is None or s["width"] is None or s["height"] is None:
                continue
            if s["is_bg"]:
                continue
            tol = OVERLAP_TOL_INCH * EMU_PER_INCH
            if s["right"] > slide_w + tol:
                amt = (s["right"] - slide_w) / EMU_PER_INCH
                msg = f"SHAPE_{s['idx']} overflow_right: shape 超出右边界 {amt:.2f} inch (type={s['type']}, text='{s['text'][:40]}')"
                slide_severe.append("[SEVERE] " + msg)
            if s["bottom"] > slide_h + tol:
                amt = (s["bottom"] - slide_h) / EMU_PER_INCH
                msg = f"SHAPE_{s['idx']} overflow_bottom: shape 超出下边界 {amt:.2f} inch (type={s['type']}, text='{s['text'][:40]}')"
                slide_severe.append("[SEVERE] " + msg)
            if s["left"] < -tol:
                amt = (-s["left"]) / EMU_PER_INCH
                msg = f"SHAPE_{s['idx']} overflow_left: shape 超出左边界 {amt:.2f} inch (type={s['type']})"
                slide_severe.append("[SEVERE] " + msg)
            if s["top"] < -tol:
                amt = (-s["top"]) / EMU_PER_INCH
                msg = f"SHAPE_{s['idx']} overflow_top: shape 超出上边界 {amt:.2f} inch (type={s['type']})"
                slide_severe.append("[SEVERE] " + msg)

        # b. 元素重叠 (过滤容器-子元素包含关系)
        non_bg = [s for s in shapes_info if not s["is_bg"] and s["rect"] is not None]
        for a_i in range(len(non_bg)):
            for b_i in range(a_i + 1, len(non_bg)):
                a = non_bg[a_i]
                b = non_bg[b_i]
                if not rect_intersect(a["rect"], b["rect"]):
                    continue
                # 两条直线相交不算
                if a["is_line"] and b["is_line"]:
                    continue
                # 容器-子元素包含关系: 一方完全包含另一方 -> 跳过 (除非两者都有文本)
                a_in_b = rect_contains(b["rect"], a["rect"])
                b_in_a = rect_contains(a["rect"], b["rect"])
                if (a_in_b or b_in_a) and not (a["text"].strip() and b["text"].strip()):
                    continue
                # 装饰小圆点(0.1x0.1)与文本重叠通常是 list bullet, 跳过
                def is_dot(s):
                    try:
                        return (s["width"] and s["height"] and
                                s["width"]/EMU_PER_INCH <= 0.15 and s["height"]/EMU_PER_INCH <= 0.15)
                    except Exception:
                        return False
                if is_dot(a) or is_dot(b):
                    continue
                ta = a["text"][:30].replace("\n", " ")
                tb = b["text"][:30].replace("\n", " ")
                msg = f"SHAPE_{a['idx']} and SHAPE_{b['idx']} overlap: '{ta}' overlaps with '{tb}' (types: {a['type']} / {b['type']})"
                if (a["text"].strip() and b["text"].strip()) or \
                   (a["text"].strip() and b["is_line"]) or (b["text"].strip() and a["is_line"]):
                    slide_severe.append("[SEVERE] " + msg)
                else:
                    slide_warnings.append("[WARNING] " + msg)

        # c. 边距不足
        margin_issue = False
        slide_w_in = slide_w / EMU_PER_INCH
        for s in shapes_info:
            if None in (s["left"], s["top"], s["width"], s["height"]):
                continue
            if s["is_bg"]:
                continue
            # 装饰小圆点边距放宽
            is_dot = (s["width"]/EMU_PER_INCH <= 0.15 and s["height"]/EMU_PER_INCH <= 0.15)
            l_in = s["left"] / EMU_PER_INCH
            t_in = s["top"] / EMU_PER_INCH
            r_in = s["right"] / EMU_PER_INCH
            if l_in < MARGIN_INCH and not is_dot:
                msg = f"SHAPE_{s['idx']} margin_left: left={l_in:.2f}in < {MARGIN_INCH}in (type={s['type']}, text='{s['text'][:30]}')"
                slide_warnings.append("[WARNING] " + msg)
                margin_issue = True
            if t_in < MARGIN_INCH and not is_dot:
                msg = f"SHAPE_{s['idx']} margin_top: top={t_in:.2f}in < {MARGIN_INCH}in (type={s['type']}, text='{s['text'][:30]}')"
                slide_warnings.append("[WARNING] " + msg)
                margin_issue = True
            if r_in > slide_w_in - MARGIN_INCH:
                msg = f"SHAPE_{s['idx']} margin_right: right={r_in:.2f}in > slide_w-{MARGIN_INCH}in={slide_w_in-MARGIN_INCH:.2f}in (type={s['type']}, text='{s['text'][:30]}')"
                slide_warnings.append("[WARNING] " + msg)
                margin_issue = True

        # d. 文本框过窄且含长文本
        for s in shapes_info:
            if s["width"] is None:
                continue
            w_in = s["width"] / EMU_PER_INCH
            text_len = len(s["text"].strip())
            if w_in < NARROW_WIDTH_INCH and text_len >= LONG_TEXT_LEN:
                msg = f"SHAPE_{s['idx']} narrow_textbox: width={w_in:.2f}in 且文本长度={text_len} (type={s['type']}, text='{s['text'][:40]}')"
                slide_warnings.append("[WARNING] " + msg)

        # e. 标题装饰线 (LINE 位于标题下方, AI 生成痕迹)
        title_candidates = []
        for s in shapes_info:
            if s["is_bg"] or not s["text"].strip():
                continue
            if s["top"] is not None and s["top"] / EMU_PER_INCH < (slide_h / EMU_PER_INCH) * 0.4:
                title_candidates.append(s)
        for s in shapes_info:
            if not s["is_line"]:
                continue
            if None in (s["left"], s["top"], s["width"], s["height"]):
                continue
            h_in = s["height"] / EMU_PER_INCH
            w_in = s["width"] / EMU_PER_INCH
            is_horizontal = (h_in < 0.15) or (w_in > h_in * 5)
            if not is_horizontal:
                continue
            for t in title_candidates:
                if t["bottom"] is None:
                    continue
                gap = (s["top"] - t["bottom"]) / EMU_PER_INCH
                if -0.1 <= gap <= 0.6:
                    if (s["left"] < t["right"] and s["right"] > t["left"]):
                        msg = f"SHAPE_{s['idx']} decorative_line_under_title: LINE 位于标题文本下方 gap={gap:.2f}in (疑似 AI 生成装饰线, width={w_in:.2f}in)"
                        slide_severe.append("[SEVERE] " + msg)
                        break

        if slide_severe:
            for line in slide_severe:
                lines_out.append(line)
                severe_issues_total += 1
        if slide_warnings:
            for line in slide_warnings:
                lines_out.append(line)
                warnings_total += 1
        if not slide_severe and not slide_warnings:
            lines_out.append("[OK] 未检测到溢出 / 重叠 / 边距 / 文本框过窄 / 标题装饰线问题")
            ok_total += 1

        # shape 清单
        lines_out.append("  --- shapes ---")
        for s in shapes_info:
            if s["left"] is None:
                lines_out.append(f"  SHAPE_{s['idx']}: type={s['type']} (无定位信息) text='{s['text'][:40]}'")
                continue
            l_in = s["left"] / EMU_PER_INCH
            t_in = s["top"] / EMU_PER_INCH
            w_in = s["width"] / EMU_PER_INCH if s["width"] else 0
            h_in = s["height"] / EMU_PER_INCH if s["height"] else 0
            lines_out.append(
                f"  SHAPE_{s['idx']}: type={s['type']} line={s['is_line']} bg={s['is_bg']} "
                f"pos=({l_in:.2f},{t_in:.2f}) size=({w_in:.2f}x{h_in:.2f}) "
                f"text='{s['text'][:50]}'"
            )

        slide_summaries.append({
            "slide": slide_idx,
            "severe": len(slide_severe),
            "warning": len(slide_warnings),
        })

    # 汇总
    lines_out.append("\n" + "=" * 70)
    lines_out.append("汇总")
    lines_out.append("=" * 70)
    lines_out.append(f"总 slide 数: {slide_count}")
    lines_out.append(f"[SEVERE] 严重问题数: {severe_issues_total}")
    lines_out.append(f"[WARNING] 警告数: {warnings_total}")
    lines_out.append(f"[OK] 无问题 slide 数: {ok_total}")
    lines_out.append("")
    lines_out.append("每张 slide 问题计数 (severe/warning):")
    for sm in slide_summaries:
        flag = ""
        if sm["severe"] > 0:
            flag = "  <<<< 严重"
        lines_out.append(f"  SLIDE {sm['slide']:>2}: severe={sm['severe']}, warning={sm['warning']}{flag}")

    report_text = "\n".join(lines_out)
    with open(REPORT_PATH, "w", encoding="utf-8") as f:
        f.write(report_text)

    print(report_text)
    print(f"\n[REPORT SAVED]: {REPORT_PATH}")


if __name__ == "__main__":
    main()
