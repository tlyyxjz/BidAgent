# -*- coding: utf-8 -*-
"""检查提交材料中的 BidAgent 字样"""
import os
import re

DEST = r"C:\Users\Lenovo\Desktop\标小智_GOAI初赛提交"

print("=== 检查 BidAgent 字样 ===\n")

for root, dirs, files in os.walk(DEST):
    for fname in sorted(files):
        fpath = os.path.join(root, fname)
        rel = os.path.relpath(fpath, DEST)

        if fname.endswith(".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(fpath)
                total = 0
                for i, slide in enumerate(prs.slides):
                    slide_hits = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text = shape.text_frame.text
                            count = text.count("BidAgent")
                            if count > 0:
                                slide_hits.append((count, text.replace("\n", "|")[:80]))
                                total += count
                    if slide_hits:
                        for c, t in slide_hits:
                            print(f"  {rel} slide{i+1} ({c}x): {t}")
                if total > 0:
                    print(f"  >>> {rel} 总计 {total} 次\n")
            except Exception as e:
                print(f"  WARN: {fname}: {e}")
        elif fname.endswith(".md"):
            with open(fpath, "r", encoding="utf-8") as f:
                lines = f.readlines()
            total = 0
            for i, line in enumerate(lines):
                count = line.count("BidAgent")
                if count > 0:
                    total += count
                    print(f"  {rel} L{i+1} ({count}x): {line.strip()[:90]}")
            if total > 0:
                print(f"  >>> {rel} 总计 {total} 次\n")
        elif fname == "LICENSE":
            with open(fpath, "r", encoding="utf-8") as f:
                content = f.read()
            if "BidAgent" in content:
                print(f"  {rel}: 含 BidAgent")
