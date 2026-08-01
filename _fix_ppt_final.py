# -*- coding: utf-8 -*-
"""
BidAgent proposal.pptx 修复脚本
任务1：删除路线图页（索引 27，第 28 页），更新页脚 N/30→N/29，团队页页码 29→28
任务2：测试数 200→826，副标题改为"项测试通过（含 parametrize 展开）"
"""
from pptx import Presentation

PPTX_PATH = r'C:\Users\Lenovo\Desktop\BidAgent\_w2_report\proposal.pptx'

prs = Presentation(PPTX_PATH)
original_count = len(prs.slides)
print(f'原始 slide 数量: {original_count}')

# ============================================================
# 任务2：修复测试数（先做，因为不涉及 slide 删除）
# 找到包含大字"200"且同页有"项测试通过"的 slide
# ============================================================
test_slide_idx = None
for i, slide in enumerate(prs.slides):
    has_200 = False
    has_test_pass = False
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t == '200':
                has_200 = True
            if '项测试通过' in t:
                has_test_pass = True
    if has_200 and has_test_pass:
        test_slide_idx = i
        break

replaced_test = []
if test_slide_idx is not None:
    slide = prs.slides[test_slide_idx]
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                old = run.text
                if run.text.strip() == '200':
                    run.text = run.text.replace('200', '826')
                    replaced_test.append((old, run.text, 'big number'))
                elif '项测试通过' in run.text:
                    run.text = '项测试通过（含 parametrize 展开）'
                    replaced_test.append((old, run.text, 'subtitle'))
    print(f'[任务2] 测试数修复完成 (slide index {test_slide_idx})')
    for old, new, label in replaced_test:
        print(f'  {label}: {repr(old)} -> {repr(new)}')
else:
    print('[任务2] 警告：未找到测试概览页！')

# ============================================================
# 任务1：删除路线图 slide
# 找到标题包含"路线图"的 slide
# ============================================================
roadmap_idx = None
roadmap_title = None
for i, slide in enumerate(prs.slides):
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if '路线图' in t and '07' in t:
                roadmap_idx = i
                roadmap_title = t
                break
    if roadmap_idx is not None:
        break

if roadmap_idx is not None:
    print(f'[任务1] 路线图 slide 找到: index={roadmap_idx}, title={repr(roadmap_title)}')
    # 用 XML 方式删除 slide（从 sldIdLst 移除引用）
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    xml_slides.remove(slides_list[roadmap_idx])
    print(f'[任务1] 路线图 slide 已删除')
else:
    print('[任务1] 警告：未找到路线图 slide！')

# ============================================================
# 更新页脚：所有 " / 30" -> " / 29"
# ============================================================
footer_count = 0
for i, slide in enumerate(prs.slides):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if '/ 30' in run.text:
                    run.text = run.text.replace('/ 30', '/ 29')
                    footer_count += 1
print(f'[页脚] 已更新 {footer_count} 处 " / 30" -> " / 29"')

# ============================================================
# 修复团队页页码：删除路线图后，团队页 footer "29 / 29" -> "28 / 29"
# （上一步把 "/ 30" -> "/ 29" 后，团队页变成 "29 / 29"，页码需减 1）
# ============================================================
team_fixed = False
for i, slide in enumerate(prs.slides):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip() == '29 / 29':
                    run.text = '28 / 29'
                    team_fixed = True
                    print(f'[页脚] 团队页页码修复: "29 / 29" -> "28 / 29" (slide index {i})')
if not team_fixed:
    print('[页脚] 警告：未找到团队页 "29 / 29" 文本！')

# ============================================================
# 更新目录：移除路线图引用
# "路线图与团队" -> "团队"
# "P.26-28" -> "P.28"
# ============================================================
toc_changes = []
for i, slide in enumerate(prs.slides):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if run.text == '路线图与团队':
                    run.text = '团队'
                    toc_changes.append(('路线图与团队', '团队'))
                elif run.text == 'P.26-28':
                    run.text = 'P.28'
                    toc_changes.append(('P.26-28', 'P.28'))
if toc_changes:
    print(f'[目录] 更新 {len(toc_changes)} 处:')
    for old, new in toc_changes:
        print(f'  {repr(old)} -> {repr(new)}')
else:
    print('[目录] 未找到需更新的目录项')

# ============================================================
# 保存到临时文件（工作目录内），后续用命令复制到目标位置
# ============================================================
import os
TEMP_OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'proposal_fixed.pptx')
prs.save(TEMP_OUT)
final_count = len(prs.slides)
print(f'\n=== 完成 ===')
print(f'slide 数量: {original_count} -> {final_count}')
print(f'删除的 slide 标题: {repr(roadmap_title) if roadmap_title else "N/A"}')
print(f'测试数替换: {replaced_test}')
print(f'已保存到临时文件: {TEMP_OUT}')
print(f'目标路径: {PPTX_PATH}')
