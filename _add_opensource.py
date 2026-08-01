"""
P2: slide 26 开放复用页——强调开源社区贡献
在底部"开源协议"行下方再加一行 GitHub 链接 + 欢迎社区贡献
"""
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor

path = r'_w2_report\proposal.pptx'
prs = Presentation(path)
slide = prs.slides[25]  # slide 26

# 在 SHAPE_28 (top=6.50) 下方加一个文本框 (top=6.80)
left = Inches(0.50)
top = Inches(6.80)
width = Inches(9.00)
height = Inches(0.25)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "github.com/tlyyxjz/BidAgent · 欢迎 issue/PR · 社区共建招投标智能体生态"
r.font.size = Pt(10)
r.font.color.rgb = RGBColor(0x16, 0x77, 0xFF)  # 蓝色链接色
r.font.bold = True

print('slide26: added opensource community line')

prs.save(path)
print('Saved.')

# 验证
prs2 = Presentation(path)
s26 = prs2.slides[25]
print(f'slide26 last 2 shapes:')
for j, sh in enumerate(list(s26.shapes)[-2:]):
    if sh.has_text_frame:
        print(f'  SHAPE: "{sh.text_frame.text}"')
