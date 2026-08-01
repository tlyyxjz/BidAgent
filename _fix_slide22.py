from pptx import Presentation

path = r'_w2_report\proposal.pptx'
prs = Presentation(path)

# 针对 slide 22 (index 21) 做更彻底的页码修复
slide = prs.slides[21]
print(f'Slide 22 shapes:')
for j, shape in enumerate(slide.shapes):
    if shape.has_text_frame:
        print(f'  SHAPE_{j}: "{shape.text_frame.text}"')
        for pi, para in enumerate(shape.text_frame.paragraphs):
            for ri, run in enumerate(para.runs):
                print(f'    para[{pi}] run[{ri}]: "{run.text}"')
                if '29' in run.text:
                    run.text = run.text.replace('29', '28')
                    print(f'      -> FIXED to "{run.text}"')
prs.save(path)

# 验证
prs2 = Presentation(path)
for i, slide in enumerate(prs2.slides):
    for shape in slide.shapes:
        if shape.has_text_frame and '/ 29' in shape.text_frame.text:
            print(f'STILL /29 in slide {i+1}')
print('Done')
